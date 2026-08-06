from __future__ import annotations

import hashlib
import io
import math
import re
import uuid
import zipfile
from dataclasses import dataclass
from functools import lru_cache
from pathlib import Path, PurePosixPath
from typing import Any

from docx import Document
from configs.settings import Settings
from db.session import SessionLocal
from models import FileAsset, FileChunk, FileExtraction, Job
from core.security import utc_now
from services.storage import LocalEncryptedStorage
from openpyxl import load_workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pypdf import PdfReader
from sqlalchemy import delete, select

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".pptx"}
CHUNK_CHARACTERS = 1600
CHUNK_OVERLAP = 160
MAX_OFFICE_ARCHIVE_MEMBERS = 5_000
MAX_OFFICE_UNCOMPRESSED_BYTES = 256 * 1024 * 1024
MAX_OFFICE_MEMBER_BYTES = 64 * 1024 * 1024
MAX_OFFICE_COMPRESSION_RATIO = 250
MAX_PDF_PAGES = 2_000
MAX_DOCX_PARAGRAPHS = 100_000
MAX_DOCX_TABLES = 5_000
MAX_XLSX_SHEETS = 256
MAX_XLSX_CELLS = 2_000_000
MAX_PPTX_SLIDES = 2_000
MAX_EXTRACTED_BLOCKS = 20_000
MAX_EXTRACTED_CHARACTERS = 5_000_000
MAX_FILE_CHUNKS = 5_000


class FileExtractionPermanentError(RuntimeError):
    def __init__(self, code: str, message: str) -> None:
        self.code = code
        super().__init__(message)


@dataclass(frozen=True)
class ExtractedBlock:
    content: str
    locator: dict[str, Any]


def _clean_text(value: str) -> str:
    return re.sub(r"\n{3,}", "\n\n", value.replace("\x00", "")).strip()


def _resource_limit() -> FileExtractionPermanentError:
    return FileExtractionPermanentError(
        "document_resource_limit",
        "文档展开后超过安全处理上限，请拆分文件后重试",
    )


def _validate_office_archive(content: bytes) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as archive:
            members = [member for member in archive.infolist() if not member.is_dir()]
            if len(members) > MAX_OFFICE_ARCHIVE_MEMBERS:
                raise _resource_limit()
            total_uncompressed = 0
            for member in members:
                path = PurePosixPath(member.filename.replace("\\", "/"))
                if path.is_absolute() or ".." in path.parts:
                    raise FileExtractionPermanentError(
                        "unsafe_office_archive",
                        "文档压缩包包含不安全路径",
                    )
                if member.flag_bits & 0x1:
                    raise FileExtractionPermanentError("encrypted_file", "暂不支持加密 Office 文件")
                total_uncompressed += member.file_size
                if (
                    member.file_size > MAX_OFFICE_MEMBER_BYTES
                    or total_uncompressed > MAX_OFFICE_UNCOMPRESSED_BYTES
                ):
                    raise _resource_limit()
                if (
                    member.file_size >= 1024 * 1024
                    and member.file_size / max(1, member.compress_size)
                    > MAX_OFFICE_COMPRESSION_RATIO
                ):
                    raise _resource_limit()
    except zipfile.BadZipFile as exc:
        raise FileExtractionPermanentError("file_parse_failed", "Office 文件结构无效") from exc


def _parse_pdf(content: bytes) -> tuple[list[ExtractedBlock], int]:
    reader = PdfReader(io.BytesIO(content))
    if reader.is_encrypted:
        raise FileExtractionPermanentError("encrypted_file", "暂不支持加密 PDF")
    if len(reader.pages) > MAX_PDF_PAGES:
        raise _resource_limit()
    blocks = []
    for index, page in enumerate(reader.pages, start=1):
        text = _clean_text(page.extract_text() or "")
        if text:
            blocks.append(ExtractedBlock(text, {"type": "page", "page": index}))
    return blocks, len(reader.pages)


def _parse_docx(content: bytes) -> tuple[list[ExtractedBlock], int]:
    document = Document(io.BytesIO(content))
    if len(document.paragraphs) > MAX_DOCX_PARAGRAPHS or len(document.tables) > MAX_DOCX_TABLES:
        raise _resource_limit()
    blocks: list[ExtractedBlock] = []
    for index, paragraph in enumerate(document.paragraphs, start=1):
        text = _clean_text(paragraph.text)
        if text:
            blocks.append(ExtractedBlock(text, {"type": "paragraph", "paragraph": index}))
    for table_index, table in enumerate(document.tables, start=1):
        rows = []
        for row in table.rows:
            rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        text = _clean_text("\n".join(rows))
        if text:
            blocks.append(ExtractedBlock(text, {"type": "table", "table": table_index}))
    return blocks, 1


def _parse_xlsx(content: bytes) -> tuple[list[ExtractedBlock], int]:
    workbook = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    try:
        if len(workbook.worksheets) > MAX_XLSX_SHEETS:
            raise _resource_limit()
        estimated_cells = sum(
            (worksheet.max_row or 0) * (worksheet.max_column or 0)
            for worksheet in workbook.worksheets
        )
        if estimated_cells > MAX_XLSX_CELLS:
            raise _resource_limit()
        blocks: list[ExtractedBlock] = []
        for worksheet in workbook.worksheets:
            last_column = get_column_letter(worksheet.max_column or 1)
            row_buffer: list[str] = []
            start_row = 1
            last_row = 1
            for row_index, row in enumerate(worksheet.iter_rows(values_only=True), start=1):
                values = ["" if value is None else str(value) for value in row]
                if not any(value.strip() for value in values):
                    continue
                if not row_buffer:
                    start_row = row_index
                row_buffer.append(" | ".join(values))
                last_row = row_index
                if sum(len(value) for value in row_buffer) < CHUNK_CHARACTERS:
                    continue
                blocks.append(
                    ExtractedBlock(
                        _clean_text("\n".join(row_buffer)),
                        {
                            "type": "cell_range",
                            "sheet": worksheet.title,
                            "range": f"A{start_row}:{last_column}{last_row}",
                        },
                    )
                )
                row_buffer = []
            if row_buffer:
                blocks.append(
                    ExtractedBlock(
                        _clean_text("\n".join(row_buffer)),
                        {
                            "type": "cell_range",
                            "sheet": worksheet.title,
                            "range": f"A{start_row}:{last_column}{last_row}",
                        },
                    )
                )
        return blocks, len(workbook.worksheets)
    finally:
        workbook.close()


def _parse_pptx(content: bytes) -> tuple[list[ExtractedBlock], int]:
    presentation = Presentation(io.BytesIO(content))
    if len(presentation.slides) > MAX_PPTX_SLIDES:
        raise _resource_limit()
    blocks: list[ExtractedBlock] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text.strip() for cell in row.cells))
        text = _clean_text("\n".join(parts))
        if text:
            blocks.append(ExtractedBlock(text, {"type": "slide", "slide": index}))
    return blocks, len(presentation.slides)


def parse_file(name: str, content: bytes) -> tuple[list[ExtractedBlock], int, str]:
    extension = Path(name).suffix.lower()
    if extension not in SUPPORTED_EXTENSIONS:
        raise FileExtractionPermanentError(
            "unsupported_file_type",
            "当前仅支持 PDF、DOCX、XLSX 和 PPTX 内容解析",
        )
    try:
        if extension in {".docx", ".xlsx", ".pptx"}:
            _validate_office_archive(content)
        if extension == ".pdf":
            blocks, page_count = _parse_pdf(content)
            parser = "pypdf"
        elif extension == ".docx":
            blocks, page_count = _parse_docx(content)
            parser = "python-docx"
        elif extension == ".xlsx":
            blocks, page_count = _parse_xlsx(content)
            parser = "openpyxl"
        else:
            blocks, page_count = _parse_pptx(content)
            parser = "python-pptx"
    except FileExtractionPermanentError:
        raise
    except Exception as exc:
        raise FileExtractionPermanentError("file_parse_failed", f"文件内容无法解析：{exc}") from exc
    if not blocks:
        raise FileExtractionPermanentError("file_has_no_text", "文件中没有可检索文本")
    if (
        len(blocks) > MAX_EXTRACTED_BLOCKS
        or sum(len(block.content) for block in blocks) > MAX_EXTRACTED_CHARACTERS
    ):
        raise _resource_limit()
    return blocks, page_count, parser


def chunk_blocks(blocks: list[ExtractedBlock]) -> list[ExtractedBlock]:
    chunks: list[ExtractedBlock] = []
    for block in blocks:
        text = block.content
        if len(text) <= CHUNK_CHARACTERS:
            chunks.append(block)
            if len(chunks) > MAX_FILE_CHUNKS:
                raise _resource_limit()
            continue
        start = 0
        part = 1
        while start < len(text):
            end = min(len(text), start + CHUNK_CHARACTERS)
            locator = {**block.locator, "part": part}
            chunks.append(ExtractedBlock(text[start:end], locator))
            if len(chunks) > MAX_FILE_CHUNKS:
                raise _resource_limit()
            if end == len(text):
                break
            start = end - CHUNK_OVERLAP
            part += 1
    return chunks


def _test_embedding(text: str, dimension: int) -> list[float]:
    seed = hashlib.shake_256(text.encode("utf-8")).digest(dimension * 2)
    values = [
        ((seed[index * 2] << 8) + seed[index * 2 + 1]) / 32767.5 - 1 for index in range(dimension)
    ]
    norm = math.sqrt(sum(value * value for value in values)) or 1
    return [value / norm for value in values]


@lru_cache(maxsize=2)
def _embedding_model(model_name: str, cache_dir: str):
    from fastembed import TextEmbedding

    return TextEmbedding(model_name=model_name, cache_dir=cache_dir, local_files_only=True)


def embed_texts(texts: list[str], settings: Settings) -> list[list[float]]:
    if settings.app_env == "test":
        return [_test_embedding(text, settings.embedding_dimension) for text in texts]
    model = _embedding_model(settings.embedding_model, str(settings.embedding_cache_dir))
    values = [list(vector) for vector in model.embed(texts)]
    if any(len(value) != settings.embedding_dimension for value in values):
        raise RuntimeError("embedding dimension does not match database schema")
    return values


def _storage(settings: Settings) -> LocalEncryptedStorage:
    return LocalEncryptedStorage(
        settings.file_storage_root,
        current_key_version=settings.file_encryption_key_version,
        key_ring=settings.file_encryption_keys(),
    )


def run_file_extract_job(job: Job, settings: Settings) -> dict[str, Any]:
    try:
        file_id = uuid.UUID(str(job.payload_json["file_id"]))
    except (KeyError, ValueError) as exc:
        raise FileExtractionPermanentError(
            "invalid_file_job", "文件解析任务缺少有效文件标识"
        ) from exc
    with SessionLocal.begin() as db:
        file_asset = db.scalar(
            select(FileAsset).where(
                FileAsset.id == file_id,
                FileAsset.enterprise_id == job.enterprise_id,
                FileAsset.deleted_at.is_(None),
            )
        )
        extraction = db.scalar(select(FileExtraction).where(FileExtraction.file_id == file_id))
        if file_asset is None or extraction is None:
            raise FileExtractionPermanentError("file_not_found", "待解析文件不存在")
        extraction.status = "processing"
        extraction.started_at = utc_now()
        extraction.error_code = None
        extraction.error_message = None
        storage_key = file_asset.storage_key
        key_version = file_asset.encryption_key_version
        file_name = file_asset.original_name
    try:
        plaintext = _storage(settings).open_decrypted(storage_key, key_version)
        blocks, page_count, parser_name = parse_file(file_name, plaintext)
        chunks = chunk_blocks(blocks)
        embeddings = embed_texts([chunk.content for chunk in chunks], settings)
    except FileExtractionPermanentError as exc:
        with SessionLocal.begin() as db:
            extraction = db.scalar(select(FileExtraction).where(FileExtraction.file_id == file_id))
            if extraction:
                extraction.status = "failed"
                extraction.error_code = exc.code
                extraction.error_message = str(exc)[:2000]
                extraction.completed_at = utc_now()
        raise
    with SessionLocal.begin() as db:
        extraction = db.scalar(
            select(FileExtraction).where(FileExtraction.file_id == file_id).with_for_update()
        )
        if extraction is None:
            raise FileExtractionPermanentError("file_not_found", "待解析文件已删除")
        db.execute(delete(FileChunk).where(FileChunk.extraction_id == extraction.id))
        for index, (chunk, embedding) in enumerate(zip(chunks, embeddings, strict=True)):
            db.add(
                FileChunk(
                    extraction_id=extraction.id,
                    file_id=file_id,
                    chunk_index=index,
                    content=chunk.content,
                    locator_json=chunk.locator,
                    token_count=max(1, len(chunk.content) // 2),
                    embedding=embedding,
                )
            )
        extraction.status = "completed"
        extraction.parser_name = parser_name
        extraction.parser_version = "1"
        extraction.page_count = page_count
        extraction.chunk_count = len(chunks)
        extraction.character_count = sum(len(chunk.content) for chunk in chunks)
        extraction.completed_at = utc_now()
        extraction.error_code = None
        extraction.error_message = None
    return {
        "file_id": str(file_id),
        "status": "completed",
        "chunks": len(chunks),
        "page_count": page_count,
    }
